import time

import wikipedia
from PyQt5 import uic
from PyQt5.QtWidgets import QWidget
from pptx.util import Pt
import urllib.request

import dialog
import loading

from dialog import CustomDialog
from pptx import Presentation
import variables

import os
import text_parser
import string


def unicode_cheker(text):
    for i in text:

        if not (
                i.isdigit() or i in string.punctuation or i.isalpha() or i in variables.ok_symb) or i in variables.broke_symb:
            text = text.replace(i, '')
    return text


def get_temp(self, id, theme, file_path):
    print(file_path)

    variables.theme = theme
    image_count = 0
    wikipedia.set_lang("ru")
    corred = wikipedia.suggest(theme)
    search = wikipedia.search(theme)
    ind = -1
    progress = 0

    if corred != None:
        theme = corred
    if search != []:
        try:
            page = wikipedia.page(theme)
        except wikipedia.DisambiguationError:
            dlg = dialog.CustomDialog(f'По запросу "{theme}" нашлось несколько результатов, пожалуйста уточните запрос')
            dlg.exec()
            exit()


    else:
        variables.curr_theme = corred

        dlg = CustomDialog(f'К сожалению нам не удалось найти "{variables.theme}"')
        if dlg.exec():
            exit()

    template = id
    presentation = Presentation(template)

    result_dir = file_path.replace('.pptx', '')
    variables.result_dir = result_dir
    os.mkdir(result_dir)
    im_save_direc = result_dir + "/images"
    txt_save_direc = result_dir + "/text"
    os.mkdir(im_save_direc)
    os.mkdir(txt_save_direc)

    file_name = file_path[::-1][:file_path[::-1].index('/')][::-1]

    with open(txt_save_direc + "/" + f'all_inf.txt', 'w', encoding='utf-8') as handler:

        handler.write(unicode_cheker(page.content))
    titles, paragraphs = text_parser.parse(txt_save_direc + "/" + f'all_inf.txt')
    s_i_count = 0
    im_names = []

    for im in page.images:

        file_ending = im.rpartition('.')[-1]
        if not (file_ending in 'png jpg gif jpeg'):
            continue
        try:
            resource = urllib.request.urlopen(im)
        except urllib.error.HTTPError:
            dlg = dialog.CustomDialog('Сервис перегружен, попробуйте позднее')
            dlg.exec()
            exit()

        out = open(im_save_direc + "/" + f'{s_i_count}.{file_ending}', 'wb')
        out.write(resource.read())
        out.close()
        im_names.append(f'{s_i_count}.{file_ending}')
        s_i_count += 1
        time.sleep(0.5)

    colored_slides_len = len(presentation.slides)
    for slide_number in range(colored_slides_len):
        slide = presentation.slides[slide_number]

        if slide_number == 0:
            slide.shapes.title.text = search[0]
            slide.placeholders[1].text = f'{page.url}'

        elif slide_number == 1:

            slide.shapes.title.text = theme
            slide.placeholders[2].text = wikipedia.summary(theme, sentences=3)
            while True:
                try:
                    slide.placeholders[1].insert_picture(im_save_direc + "/" + im_names[image_count])
                    break
                except BaseException:
                    image_count += 1
                if s_i_count < image_count:
                    break

            image_count += 1
        else:
            print(slide_number)
            if len(paragraphs) >= slide_number:
                slide.shapes.title.text = titles[slide_number - 2]
                slide.placeholders[1].text = ''.join(paragraphs[slide_number - 2])
                for paragraph in slide.placeholders[1].text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(11)
    print(f'max was {variables.max_slides}')
    for slide_number, slide_text in enumerate(paragraphs[len(presentation.slides) - 1:]):
        print(slide_number, slide)
        if slide_text == [] and titles[slide_number + colored_slides_len - 1] == []:
            break
        if slide_number >= variables.max_slides - colored_slides_len and variables.max_slides - colored_slides_len != 0:
            break
        if 'Ссыл' in titles[slide_number + colored_slides_len - 1]:
            ind = slide_number + colored_slides_len - 1
            continue
        slide.shapes.title.text = titles[slide_number + colored_slides_len - 1]

        slide = presentation.slides.add_slide(presentation.slide_layouts[8])
        slide.placeholders[2].text = ''.join(slide_text)
        if s_i_count > image_count:
            slide.placeholders[1].insert_picture(im_save_direc + "/" + im_names[image_count])
        for paragraph in slide.placeholders[2].text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(11)
        image_count += 1
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = 'Источники'

    slide.placeholders[1].text = ''.join(paragraphs[ind])

    presentation.save(result_dir + "/" + file_name)
    variables.total = result_dir + "/" + file_name
